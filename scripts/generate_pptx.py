#!/usr/bin/env python3
"""
尊嚴農業網站架構簡報產生器
使用 python-pptx 產生 PowerPoint 簡報檔案

執行方式：python3 scripts/generate_pptx.py
產出路徑：docs/website-architecture-presentation.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── 品牌色常數 ──
BRAND_GREEN = RGBColor(0x5C, 0x86, 0x07)
BRAND_GREEN_DARK = RGBColor(0x4A, 0x6B, 0x05)
BRAND_BG = RGBColor(0xE3, 0xE9, 0xD8)
BRAND_GOLD = RGBColor(0xD4, 0xA0, 0x17)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x1F, 0x29, 0x37)
GRAY_TEXT = RGBColor(0x6B, 0x72, 0x80)
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)

# 分眾色
SAGE_GREEN = RGBColor(0x87, 0xAE, 0x73)
SLATE_BLUE = RGBColor(0x5B, 0x7F, 0xA6)
TERRACOTTA = RGBColor(0xC1, 0x7F, 0x5E)
HARVEST_GOLD = RGBColor(0xD4, 0xA0, 0x17)

# 危害控制色
EMERALD = RGBColor(0x10, 0xB9, 0x81)
SKY_BLUE = RGBColor(0x0E, 0xA5, 0xE9)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
ROSE = RGBColor(0xF4, 0x3F, 0x5E)

OUTPUT_PATH = os.path.join(os.path.dirname(__file__), '..', 'docs', 'website-architecture-presentation.pptx')


def add_brand_bar(slide, color=BRAND_GREEN):
    """在投影片頂部加入品牌色條"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.15))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_footer(slide, text="尊嚴農業 — 農業健康與安全知識樞紐"):
    """在投影片底部加入頁腳"""
    left = Inches(0.5)
    top = Inches(7.0)
    width = Inches(12)
    height = Inches(0.4)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = GRAY_TEXT
    p.alignment = PP_ALIGN.RIGHT


def set_slide_bg(slide, color=WHITE):
    """設定投影片背景色"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_slide(prs, title, subtitle, bg_color=BRAND_GREEN):
    """建立標題投影片（封面或段落分隔頁）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, bg_color)

    # 標題
    left = Inches(1)
    top = Inches(2.2)
    width = Inches(11)
    height = Inches(2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # 副標題
    top2 = Inches(4.5)
    txBox2 = slide.shapes.add_textbox(left, top2, width, Inches(1.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(0xDD, 0xED, 0xBB)
    p2.alignment = PP_ALIGN.CENTER

    return slide


def add_content_slide(prs, title, bullets, note=""):
    """建立內容投影片（標題 + 條列）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, WHITE)
    add_brand_bar(slide)

    # 標題
    left = Inches(0.7)
    top = Inches(0.4)
    width = Inches(12)
    height = Inches(0.9)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = BRAND_GREEN_DARK
    p.alignment = PP_ALIGN.LEFT

    # 標題下方裝飾線
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.3), Inches(2), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = BRAND_GOLD
    line.line.fill.background()

    # 內容
    content_top = Inches(1.6)
    content_height = Inches(5.0)
    txBox2 = slide.shapes.add_textbox(Inches(0.9), content_top, Inches(11.5), content_height)
    tf2 = txBox2.text_frame
    tf2.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()

        # 支援巢狀層級（用 tuple: (level, text)）
        if isinstance(bullet, tuple):
            level, text = bullet
            p.level = level
            p.text = text
        else:
            p.text = bullet
            p.level = 0

        p.font.size = Pt(16)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(8)

    if note:
        p_note = tf2.add_paragraph()
        p_note.text = note
        p_note.font.size = Pt(12)
        p_note.font.italic = True
        p_note.font.color.rgb = GRAY_TEXT
        p_note.space_before = Pt(16)

    add_footer(slide)
    return slide


def add_two_column_slide(prs, title, left_title, left_bullets, right_title, right_bullets):
    """建立雙欄投影片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_brand_bar(slide)

    # 標題
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(12), Inches(0.9))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = BRAND_GREEN_DARK

    # 裝飾線
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.3), Inches(2), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = BRAND_GOLD
    line.line.fill.background()

    # 左欄
    _add_column(slide, Inches(0.7), left_title, left_bullets)
    # 右欄
    _add_column(slide, Inches(6.8), right_title, right_bullets)

    add_footer(slide)
    return slide


def _add_column(slide, left, col_title, bullets):
    """輔助：加入單欄內容"""
    # 欄標題
    txBox = slide.shapes.add_textbox(left, Inches(1.6), Inches(5.5), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = col_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = BRAND_GREEN

    # 欄內容
    txBox2 = slide.shapes.add_textbox(left, Inches(2.3), Inches(5.5), Inches(4.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = bullet
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(6)


def add_color_block(slide, left, top, width, height, color, label, sublabel=""):
    """加入色塊標示"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if sublabel:
        p2 = tf.add_paragraph()
        p2.text = sublabel
        p2.font.size = Pt(10)
        p2.font.color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
        p2.alignment = PP_ALIGN.CENTER


def add_table_slide(prs, title, headers, rows):
    """建立表格投影片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_brand_bar(slide)

    # 標題
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(12), Inches(0.9))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = BRAND_GREEN_DARK

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.3), Inches(2), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = BRAND_GOLD
    line.line.fill.background()

    # 表格
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table_left = Inches(0.7)
    table_top = Inches(1.6)
    table_width = Inches(12)
    table_height = Inches(0.4) * num_rows

    table_shape = slide.shapes.add_table(num_rows, num_cols, table_left, table_top, table_width, table_height)
    table = table_shape.table

    # 設定欄寬（平均分配）
    col_width = int(Inches(12) / num_cols)
    for i in range(num_cols):
        table.columns[i].width = col_width

    # 表頭
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(13)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = BRAND_GREEN

    # 表格內容
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = val
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.font.color.rgb = DARK_TEXT
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BG
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE

    add_footer(slide)
    return slide


# ═══════════════════════════════════════════════
# 主程式：建立所有投影片
# ═══════════════════════════════════════════════

def build_presentation():
    prs = Presentation()
    # 設定為 16:9 寬螢幕
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ──────────────────────────────────────────
    # 投影片 1：封面
    # ──────────────────────────────────────────
    slide = add_title_slide(
        prs,
        "尊嚴農業\n網站架構與內容說明",
        "農業尊嚴：共建雲嘉農業工作者職業健康社會安全網\n"
        "國立中正大學 × 臺大醫院雲林分院｜嘉義縣民雄鄉",
        bg_color=BRAND_GREEN_DARK
    )
    # 底部標記
    txBox = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Knowledge Hub for Agricultural Health & Safety"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0xAA, 0xCC, 0x77)
    p.alignment = PP_ALIGN.CENTER

    # ──────────────────────────────────────────
    # 投影片 2：問題意識
    # ──────────────────────────────────────────
    add_content_slide(
        prs,
        "為什麼需要這個網站？",
        [
            "台灣農業面臨的嚴峻挑戰：",
            (1, "農保被保險人平均年齡 67.3 歲 — 農業勞動力高度老化"),
            (1, "農業致死率為全行業平均的 5 倍（CDC/NIOSH, 2024）"),
            (1, "全球每年約 38.5 萬例急性農藥中毒（ILO）"),
            (1, "農民職業災害保險投保率僅約 30%（勞保局統計）"),
            (1, "農業致死者中 56% 年齡 ≥ 55 歲（CDC/NIOSH）"),
            "",
            "傳統農業資訊網站的問題：",
            (1, "政令宣導式條列網頁，缺乏分眾設計"),
            (1, "資訊被動、單向，無互動工具"),
            (1, "缺乏國際數據佐證與學術嚴謹度"),
        ],
        note="資料來源：CDC/NIOSH (2024), ILO (2023), WHO (2022), 台灣農保統計, 勞保局"
    )

    # ──────────────────────────────────────────
    # 投影片 3：矩陣式資訊架構
    # ──────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_brand_bar(slide)

    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(12), Inches(0.9))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "解決方案：矩陣式資訊架構"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = BRAND_GREEN_DARK

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.3), Inches(2), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = BRAND_GOLD
    line.line.fill.background()

    # 說明文字
    txBox2 = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(5), Inches(1))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "首頁＝分流大廳（The Lobby）"
    p2.font.size = Pt(16)
    p2.font.bold = True
    p2.font.color.rgb = DARK_TEXT
    p3 = tf2.add_paragraph()
    p3.text = "降低認知負荷，精準內容策展"
    p3.font.size = Pt(14)
    p3.font.color.rgb = GRAY_TEXT

    # 矩陣圖 - 橫軸標籤（左側）
    audiences = [
        ("新手務農 Start", SAGE_GREEN),
        ("青農創業 Pro", SLATE_BLUE),
        ("資深農友 Senior", TERRACOTTA),
        ("一般公眾 Public", HARVEST_GOLD),
    ]
    y_start = Inches(2.8)
    for i, (label, color) in enumerate(audiences):
        add_color_block(slide,
                        Inches(0.7), y_start + Inches(i * 1.0),
                        Inches(2.5), Inches(0.8),
                        color, label)

    # 矩陣圖 - 縱軸標籤（上方）
    modules = ["01\n關於", "02\n職安", "03\n健康", "04\n保險", "05\n培力", "06\n倡議", "07\n研究"]
    x_start = Inches(3.5)
    for i, mod in enumerate(modules):
        add_color_block(slide,
                        x_start + Inches(i * 1.35), Inches(2.0),
                        Inches(1.2), Inches(0.7),
                        BRAND_GREEN, mod)

    # 矩陣交叉區域
    matrix_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        x_start, y_start,
        Inches(1.35 * 7), Inches(4.0)
    )
    matrix_bg.fill.solid()
    matrix_bg.fill.fore_color.rgb = BRAND_BG
    matrix_bg.line.color.rgb = BRAND_GREEN
    matrix_bg.line.width = Pt(1)

    matrix_text = matrix_bg.text_frame
    matrix_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = matrix_text.paragraphs[0]
    p.text = "\n\n依使用者身分\n動態撈取對應縱軸模組的精選內容"
    p.font.size = Pt(14)
    p.font.color.rgb = BRAND_GREEN_DARK
    p.font.italic = True

    # 標註
    txBox3 = slide.shapes.add_textbox(Inches(0.7), Inches(6.8), Inches(12), Inches(0.4))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "橫軸 = 分眾入口（依生命週期）　　縱軸 = 結構化資料庫（依知識領域）"
    p3.font.size = Pt(12)
    p3.font.color.rgb = GRAY_TEXT
    p3.alignment = PP_ALIGN.CENTER

    add_footer(slide)

    # ──────────────────────────────────────────
    # 投影片 4：四大分眾入口總覽
    # ──────────────────────────────────────────
    add_table_slide(
        prs,
        "橫軸：四大分眾入口",
        ["分眾", "定位", "主色調", "介面模式", "核心功能"],
        [
            ["新手務農", "防禦包", "鼠尾草綠 #87AE73", "檢核清單 + 進度條", "務農風險快篩表 → PPE 型錄"],
            ["青農創業", "升級包", "板岩藍 #5B7FA6", "儀表板 + 比較圖表", "投資報酬率 + 退休金試算"],
            ["資深農友", "樂活包", "陶土橘 #C17F5E", "Mobile-First + App", "找醫師・算津貼・看影片"],
            ["一般公眾", "認同包", "麥穗黃 #D4A017", "敘事性滾動", "攝影紀實 + Podcast + 倡議"],
        ]
    )

    # ──────────────────────────────────────────
    # 投影片 5：新手務農
    # ──────────────────────────────────────────
    slide = add_content_slide(
        prs,
        "新手務農 — 防禦包",
        [
            "開場震撼：「你知道農業是全台最危險的職業之一嗎？」",
            "",
            "設計策略：",
            (1, "鼠尾草綠（Sage Green）— 傳遞穩健、安全感"),
            (1, "左側固定導覽步驟，右側動態表單"),
            (1, "首屏強制顯示「務農風險快篩表」"),
            "",
            "關鍵數據：",
            (1, "5× — 農業致死率是全行業平均的 5 倍（CDC/NIOSH）"),
            (1, "29% — 農業傷害由跌倒造成（CDC/NIOSH）"),
            (1, "38.5 萬 — 全球每年農藥急性中毒人數（ILO）"),
        ]
    )

    # ──────────────────────────────────────────
    # 投影片 6：青農創業
    # ──────────────────────────────────────────
    add_content_slide(
        prs,
        "青農創業 — 升級包",
        [
            "開場共鳴：「你已經在田裡了。十年後你的身體還撐得住嗎？」",
            "",
            "設計策略：",
            (1, "板岩藍（Slate Blue）— 科技感、數據導向"),
            (1, "左側：省工農機投資報酬率量化圖表"),
            (1, "右側：農民退休儲金 vs. 勞保高階比較試算器"),
            "",
            "關鍵數據：",
            (1, "67.3 歲 — 農保被保險人平均年齡（台灣農保統計）"),
            (1, "1/3 ～ 1/2 — 農機政府補助比例（農業部農糧署）"),
            (1, "30% — 農民職業災害保險投保率（勞保局統計）"),
        ]
    )

    # ──────────────────────────────────────────
    # 投影片 7：資深農友
    # ──────────────────────────────────────────
    add_content_slide(
        prs,
        "資深農友 — 樂活包",
        [
            "開場糾正：「腰痛不是老了，是工作在傷害你。」",
            "",
            "設計策略：",
            (1, "陶土橘（Terracotta）— 溫暖、高對比"),
            (1, "Mobile-First + App-like Interface"),
            (1, "最小字體 18px、最小按鈕高度 48px 觸控熱區"),
            "",
            "底部固定導覽（三大入口）：",
            (1, "找醫師（地圖）"),
            (1, "算津貼（計算機）"),
            (1, "看影片（微學習）"),
            "",
            "個案故事：曾女士 — 民雄農友，40 年務農後成功申請職傷給付",
        ]
    )

    # ──────────────────────────────────────────
    # 投影片 8：一般公眾
    # ──────────────────────────────────────────
    add_content_slide(
        prs,
        "一般公眾 — 認同包",
        [
            "開場連結：「你在超市拿起的平價蔬果，種它的人正忍著腰痛。」",
            "",
            "設計策略：",
            (1, "麥穗黃（Harvest Gold）— 農鄉共感"),
            (1, "敘事性滾動（Scrollytelling）"),
            (1, "滿版高畫質雲嘉農鄉紀實攝影"),
            "",
            "特色元素：",
            (1, "雲嘉農業職災熱力圖"),
            (1, "高齡化數據資訊圖表（Infographics）"),
            (1, "Podcast 播放器 — 農民訪談錄音"),
            "",
            "引言：陳宗延醫師（臺大醫院雲林分院）",
            (1, "「農民沒有慣老闆，但他們有兩個更殘酷的超級老闆：天氣和市場。」"),
        ]
    )

    # ──────────────────────────────────────────
    # 投影片 9：七大資料庫總覽
    # ──────────────────────────────────────────
    add_table_slide(
        prs,
        "縱軸：七大資料庫主選單",
        ["編號", "模組名稱", "介面類型", "核心功能"],
        [
            ["01", "關於計畫", "靜態資訊 + 多媒體", "團隊卡片 + 互動地圖"],
            ["02", "職業安全", "網格卡片 + 互動工具", "NIOSH 控制層級 + 紅綠燈風險評估"],
            ["03", "健康促進", "地理資訊 + 影音", "互動人體圖 + 職醫地圖"],
            ["04", "經濟與保險", "圖解 + 試算工具", "保險試算器 + 法規懶人包"],
            ["05", "培力與實踐", "輕量 LMS + 行事曆", "影片網格 + 工作坊報名"],
            ["06", "尊嚴農業倡議", "單頁式深度專題", "垂直時間軸 + 政策連署"],
            ["07", "研究成果", "學術資料庫", "APA 7 文獻 + Podcast 嵌入"],
        ]
    )

    # ──────────────────────────────────────────
    # 投影片 10：職業安全模組
    # ──────────────────────────────────────────
    add_content_slide(
        prs,
        "02 職業安全模組 — 深度解析",
        [
            "頁面結構：",
            (1, "Hero 區塊（標題 + 數據呼叫盒）"),
            (1, "Tabs 切換：省工農機｜個人防護具 PPE"),
            (1, "Grid Cards 數位型錄（含 PDF 圖卡下載）"),
            (1, "風險評估工具區"),
            "",
            "互動工具：",
            (1, "紅綠燈快篩系統 — 高（紅）/ 中（橘）/ 低（綠）"),
            (1, "NIOSH 五層安全控制層級動態圖"),
            (1, "四大危害深度檢核：生物・物理・化學・人因"),
            "",
            "子頁面：溫室專區｜省工農機具｜個人防護具",
            "",
            "數據佐證：CDC/NIOSH 18.6/10 萬 FTE 致死率、BLS 拖拉機翻覆為首因",
        ]
    )

    # ──────────────────────────────────────────
    # 投影片 11：健康促進 & 經濟保險
    # ──────────────────────────────────────────
    add_two_column_slide(
        prs,
        "03 健康促進 & 04 經濟與保險",
        "03 健康促進",
        [
            "互動式人體圖",
            "  → 點擊腰/膝/肩",
            "  → Modal 播放復健 GIF/影片",
            "",
            "職醫地圖（Google Maps API）",
            "  → 勞動部認可診治網絡醫院",
            "  → 一鍵導航",
            "",
            "即時氣象高溫警報",
            "急救紅卡（可下載列印）",
        ],
        "04 經濟與保險",
        [
            "互動試算工具",
            "  → 輸入年齡/年資/傷病類型",
            "  → 輸出可申請補助或理賠估算",
            "",
            "涵蓋險種：",
            "  → 職災保險",
            "  → 農保・老農津貼",
            "  → 農業保險",
            "",
            "法規→步驟懶人包（流程圖）",
            "退休金試算表",
        ]
    )

    # ──────────────────────────────────────────
    # 投影片 12：培力/倡議/研究
    # ──────────────────────────────────────────
    add_content_slide(
        prs,
        "05 培力與實踐 ｜ 06 倡議 ｜ 07 研究成果",
        [
            "05 培力與實踐（農學堂）：",
            (1, "影片網格（Video Grid）— 每支標示時長與難易度"),
            (1, "工作坊清單 — 雲嘉近期線下課程 + 內建報名表單"),
            "",
            "06 尊嚴農業倡議：",
            (1, "垂直時間軸（Vertical Timeline）— 串接倡議行動歷程"),
            (1, "焦點團體會議紀錄"),
            (1, "頁面底部：政策連署 + 社群分享按鈕"),
            "",
            "07 研究成果：",
            (1, "嚴格 APA 第 7 版格式排版"),
            (1, "每筆文獻附 [PDF 下載, X.X MB] 按鈕"),
            (1, "整合 Spotify / Apple Podcast 嵌入代碼"),
            (1, "子頁面：下載專區｜圖文資訊｜Podcast｜相關資源｜研究出版"),
        ]
    )

    # ──────────────────────────────────────────
    # 投影片 13：設計系統
    # ──────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_brand_bar(slide)

    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(12), Inches(0.9))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "設計系統 — 色彩與字型"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = BRAND_GREEN_DARK

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.3), Inches(2), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = BRAND_GOLD
    line.line.fill.background()

    # 品牌色色塊
    colors_data = [
        ("尊嚴綠\n#5C8607", BRAND_GREEN),
        ("深尊嚴綠\n#4A6B05", BRAND_GREEN_DARK),
        ("稻金黃\n#D4A017", BRAND_GOLD),
        ("背景色\n#E3E9D8", BRAND_BG),
    ]
    for i, (label, color) in enumerate(colors_data):
        text_color = DARK_TEXT if color == BRAND_BG else WHITE
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.7 + i * 2.2), Inches(1.7),
            Inches(2.0), Inches(1.0)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        tf = shape.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = text_color

    # 分眾色色塊
    audience_colors = [
        ("新手務農\n鼠尾草綠", SAGE_GREEN),
        ("青農創業\n板岩藍", SLATE_BLUE),
        ("資深農友\n陶土橘", TERRACOTTA),
        ("一般公眾\n麥穗黃", HARVEST_GOLD),
    ]
    for i, (label, color) in enumerate(audience_colors):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.7 + i * 2.2), Inches(3.0),
            Inches(2.0), Inches(1.0)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        tf = shape.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = WHITE

    # NIOSH 控制層級色
    niosh_colors = [
        ("L1 消除", EMERALD),
        ("L2 取代", SKY_BLUE),
        ("L3 工程", BLUE),
        ("L4 行政", AMBER),
        ("L5 PPE", ROSE),
    ]
    for i, (label, color) in enumerate(niosh_colors):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.7 + i * 1.8), Inches(4.3),
            Inches(1.6), Inches(0.7)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        tf = shape.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = WHITE

    # 字型說明
    txBox3 = slide.shapes.add_textbox(Inches(0.7), Inches(5.3), Inches(12), Inches(1.5))
    tf3 = txBox3.text_frame
    tf3.word_wrap = True
    items = [
        "字型系統：Noto Sans TC（中文）｜Lato（英數）｜Lora（內文裝飾）",
        "圖示庫：Font Awesome 6.4.x",
        "視覺風格：功能性柔和色系（Functional Pastels）",
        "CSS 設計 Token：所有色碼透過 :root 變數統一管理",
    ]
    for i, item in enumerate(items):
        if i == 0:
            p = tf3.paragraphs[0]
        else:
            p = tf3.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(4)

    add_footer(slide)

    # ──────────────────────────────────────────
    # 投影片 14：技術架構
    # ──────────────────────────────────────────
    add_content_slide(
        prs,
        "技術架構",
        [
            "平台基礎：",
            (1, "WordPress CMS + Astra 主題"),
            (1, "HTML 區塊直貼 WordPress 程式碼編輯器（無需額外部署）"),
            "",
            "前端技術：",
            (1, "純 HTML/CSS — 無需 build 工具"),
            (1, "CSS Custom Properties（設計 Token）統一管理全站樣式"),
            (1, "CSS 前綴隔離策略（如 gh- / mach- / ppe-）防止跨頁衝突"),
            (1, "Google Fonts CDN + Font Awesome 6.4.x CDN"),
            "",
            "後端功能：",
            (1, "PHP Shortcodes — 新聞迴圈、Podcast 網格、天氣小工具、FB 動態牆"),
            (1, "WordPress 文章分類系統（階層式 slug）"),
            (1, "WP_Query 依分類動態撈取內容"),
            "",
            "WordPress 嵌入規範：禁止 body/html/全域重置等覆蓋主題樣式的寫法",
        ]
    )

    # ──────────────────────────────────────────
    # 投影片 15：程式碼結構
    # ──────────────────────────────────────────
    add_content_slide(
        prs,
        "程式碼倉庫結構",
        [
            "farmerwordpress/",
            (1, "docs/          — 架構藍圖、設計規格文件"),
            (1, "css/           — global.css 全站設計 Token"),
            (1, "pages/         — 各頁面 HTML 區塊"),
            (1, "  ├── About/            — 01 關於計畫"),
            (1, "  ├── occupational-safety/ — 02 職業安全（含子頁）"),
            (1, "  ├── beginner-farmers/  — 新手務農"),
            (1, "  ├── young-farmers/     — 青農"),
            (1, "  ├── Public/            — 一般民眾"),
            (1, "  └── shared/            — 共用元件（sidebar）"),
            (1, "shortcodes/    — WordPress PHP 短代碼"),
            (1, "scripts/       — 部署與匯入腳本"),
            (1, "imports/       — 批次匯入 HTML 暫存區"),
            "",
            "PHP Shortcodes：",
            (1, "新聞迴圈（my_news_loop）｜Podcast 網格（knowledge_podcast_grid）"),
            (1, "天氣小工具（weather_widget）｜FB 動態牆（fb_feed_widget）"),
            (1, "首頁亮點輪播（home_highlight_slider）｜雜誌式新聞版面（news_magazine_layout）"),
        ]
    )

    # ──────────────────────────────────────────
    # 投影片 16：學術嚴謹度
    # ──────────────────────────────────────────
    add_content_slide(
        prs,
        "學術嚴謹度與數據治理",
        [
            "引用規範：",
            (1, "所有文獻嚴格遵守 APA 第 7 版格式"),
            (1, "系統性排除中國大陸來源數據"),
            (1, "以台灣本土實證 + 日韓澳相似高齡農業背景研究為學術基石"),
            "",
            "國際權威機構引用清單：",
            (1, "CDC/NIOSH — 農業安全統計、控制層級、抬舉方程式"),
            (1, "OSHA — 農業危害與控制、熱傷害預防、PPE 標準"),
            (1, "BLS — 致命職業傷害普查（CFOI）"),
            (1, "WHO — 農藥安全（2022）"),
            (1, "ILO — 農業安全與健康（2023）"),
            "",
            "台灣本土數據：",
            (1, "農保統計（平均年齡 67.3 歲）｜勞保局（投保率 30%）"),
            (1, "農業部農糧署（農機補助）｜勞動部職安署"),
        ]
    )

    # ──────────────────────────────────────────
    # 投影片 17：開發進度路線圖
    # ──────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_brand_bar(slide)

    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(12), Inches(0.9))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "開發進度路線圖"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = BRAND_GREEN_DARK

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.3), Inches(2), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = BRAND_GOLD
    line.line.fill.background()

    # Phase 色塊
    phases = [
        ("Phase 1", "已完成", "首頁：身分導航 + 三大核心主題", RGBColor(0x10, 0xB9, 0x81), "DONE"),
        ("Phase 2", "進行中", "職業安全頁面重建\n（主頁 + 溫室 + 省工農機 + PPE）", RGBColor(0x3B, 0x82, 0xF6), "WIP"),
        ("Phase 3", "待啟動", "健康促進\n（人體圖 + 職醫地圖）\n經濟試算器", RGBColor(0xF5, 0x9E, 0x0B), "NEXT"),
        ("Phase 4", "規劃中", "分眾入口頁面\n農學堂 LMS\n研究成果資料庫", RGBColor(0x6B, 0x72, 0x80), "PLAN"),
    ]

    for i, (phase, status, desc, color, badge) in enumerate(phases):
        x = Inches(0.7 + i * 3.1)
        y = Inches(1.8)

        # 色塊
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.8), Inches(4.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Phase 標題
        p1 = tf.paragraphs[0]
        p1.text = f"{phase}"
        p1.font.size = Pt(20)
        p1.font.bold = True
        p1.font.color.rgb = WHITE

        # 狀態
        p2 = tf.add_paragraph()
        p2.text = status
        p2.font.size = Pt(14)
        p2.font.color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
        p2.alignment = PP_ALIGN.CENTER

        # 分隔
        p3 = tf.add_paragraph()
        p3.text = "─────"
        p3.font.size = Pt(10)
        p3.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        p3.alignment = PP_ALIGN.CENTER

        # 內容
        p4 = tf.add_paragraph()
        p4.text = desc
        p4.font.size = Pt(13)
        p4.font.color.rgb = WHITE
        p4.alignment = PP_ALIGN.CENTER

    add_footer(slide)

    # ──────────────────────────────────────────
    # 投影片 18：社會影響力
    # ──────────────────────────────────────────
    add_content_slide(
        prs,
        "社會影響力與政策連結",
        [
            "覆蓋範圍：",
            (1, "嘉義縣民雄鄉為核心實踐場域"),
            (1, "雲嘉地區農業工作者為主要服務對象"),
            (1, "全球 8.79 億農業勞動者的共通議題（ILO）"),
            "",
            "政策對接：",
            (1, "勞動部 — 農民職業災害保險制度"),
            (1, "農業部 — 農機補助、農民退休儲金"),
            (1, "衛福部 — 職業傷病診治網絡"),
            "",
            "跨國比較研究：",
            (1, "日本 — 高齡農業社會安全網經驗"),
            (1, "韓國 — 農業職業健康政策"),
            (1, "澳洲 — 農場安全標準"),
            "",
            "知識樞紐定位：連結學術研究、政策推動、田間實務",
        ]
    )

    # ──────────────────────────────────────────
    # 投影片 19：結語
    # ──────────────────────────────────────────
    add_title_slide(
        prs,
        "讓每一位農業工作者\n都能有尊嚴地工作",
        "農業尊嚴：共建雲嘉農業工作者職業健康社會安全網\n\n感謝聆聽",
        bg_color=BRAND_GREEN
    )

    # 儲存
    output = os.path.abspath(OUTPUT_PATH)
    os.makedirs(os.path.dirname(output), exist_ok=True)
    prs.save(output)
    print(f"簡報已產生：{output}")
    print(f"投影片數量：{len(prs.slides)} 頁")


if __name__ == '__main__':
    build_presentation()
